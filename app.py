from __future__ import annotations

import argparse
import base64
import hashlib
import json
import locale
import math
import os
import re
import subprocess
import traceback
import unicodedata
from copy import deepcopy
from pathlib import Path
from typing import Any

from openai import OpenAI
from PIL import Image
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

BASE_DIR = Path(__file__).resolve().parent
TEMPLATE_DIR = BASE_DIR / "template"
PROMPTS_DIR = BASE_DIR / "prompts"
OUTPUT_DIR = BASE_DIR / "output"
IMAGES_DIR = OUTPUT_DIR / "images"
GOOGLE_SAFE_MODE = True
LANDSCAPE_IMAGE_SAFE_MARGIN_RATIO = 0.08
LANDSCAPE_IMAGE_LEFT_SAFE_MARGIN_RATIO = 0.14

DEFAULT_PROMPT_FILE = PROMPTS_DIR / "lecture_prompt.txt"
DEFAULT_CURRICULUM_FILE = PROMPTS_DIR / "curriculum.md"
DIAGRAM_FILL_RGB = RGBColor(157, 195, 230)  # #9DC3E6
DIAGRAM_LINE_RGB = RGBColor(91, 155, 213)
DEFAULT_ENV_FILE = BASE_DIR / ".env"


JSON_FORMAT_GUIDE = """
Return JSON only.

{
  "deck_title": "string",
  "section_label": "string",
  "slides": [
    {
      "slide_number": 1,
      "slide_role": "objective | content | practice_problem | practice_answer | summary",
      "source_slide": 21,
      "title": "string",
      "why": "string",
      "bullets": [
        "string",
        {"text": "string", "level": 1, "bullet": true}
      ],
      "example": "string",
      "practice_prompt": "string",
      "transition": "string",
      "image_prompt": "string",
      "visual_type": "bullets | table | diagram",
      "table": {
        "headers": ["string", "string"],
        "rows": [["string", "string"]]
      },
      "diagram": {
        "diagram_type": "process | comparison | hierarchy | cycle | relationship",
        "layout_direction": "horizontal | vertical | radial",
        "title": "string",
        "nodes": ["string", "string"],
        "links": [["string", "string"]],
        "notes": "string"
      }
    }
  ]
}

Rules:
- Keep exactly 20 slides unless the prompt explicitly requests another count.
- Slide 1 must use slide_role "objective".
- Slide 18 must use slide_role "practice_problem".
- Slide 19 must use slide_role "practice_answer".
- Slide 20 must use slide_role "summary".
- Use the provided content template slide as the main source_slide unless there is a strong reason not to.
- The template structure is simple: left-top bullet text area and right-bottom image area.
- Every slide should choose the clearest presentation form among bullets, table, or diagram.
- "bullets" may contain plain strings or objects with fields {text, level, bullet}.
- Use level 1 or level 2 bullet objects when a hierarchical explanation is clearer.
- Use the "example" field only when an inline example is truly useful. It will be rendered in the slide body as an `ex>` line.
- Set "visual_type" to "bullets" when normal bullet explanation is enough.
- Set "visual_type" to "table" when the content is clearer as rows and columns.
- Set "visual_type" to "diagram" when the content is clearer as process flow, comparison, hierarchy, cycle, or relationship.
- When "visual_type" is "table", fill the "table" object with practical headers and rows.
- When "visual_type" is "diagram", fill the "diagram" object with structured nodes and links.
- When "visual_type" is "diagram", also choose "layout_direction" based on whether the flow reads better horizontally, vertically, or radially.
- Every bullet must be a presentation-ready sentence of about 18 to 35 Korean characters when possible.
- Avoid essay-like long sentences.
- Avoid keyword-only short fragments.
- Prefer one idea per bullet and keep the wording easy to say aloud in class.
- Prefer concise nominal endings such as "~?", "~?", "~?", "~??", "~??".
- Avoid sentence endings like "~?", "~???", "~??" in slide bullets.
- Put the image prompt in "image_prompt" when a visual aid is useful.
- Make image prompts favor a wide horizontal composition for a slide image frame.
- Keep the main subject centered with safe margins on all sides.
- Avoid important objects touching the edges and avoid text near the image border.
""".strip()


def load_env_file(path: Path) -> None:
    if not path.exists():
        return

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        normalized_key = key.strip().lstrip("\ufeff")
        os.environ.setdefault(normalized_key, value.strip())


def env_flag(name: str, default: bool) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a lecture deck from an LLM prompt file and a PPT template."
    )
    parser.add_argument(
        "--prompt-file",
        default=str(DEFAULT_PROMPT_FILE),
        help="Prompt text file path. Default: ./prompts/lecture_prompt.txt",
    )
    parser.add_argument(
        "--template",
        help="Template .pptx path. Default: first .pptx in ./template",
    )
    parser.add_argument(
        "--curriculum-file",
        help="Curriculum markdown/text file path for multi-session generation. Default: ./prompts/curriculum.md",
    )
    parser.add_argument(
        "--lecture",
        help="Comma-separated lecture/session numbers to run. Example: 1 or 1,3",
    )
    parser.add_argument(
        "--page",
        help="Comma-separated slide numbers to render within selected lecture(s). Requires --lecture. Example: 1 or 1,3",
    )
    parser.add_argument(
        "--output",
        help="Output .pptx path. Default: ./output/generated_<template>.pptx",
    )
    parser.add_argument(
        "--json-output",
        help="Where to save the structured slide JSON. Default: ./output/generated_slide_plan.json",
    )
    parser.add_argument(
        "--analysis-output",
        help="Where to save the template analysis JSON. Default: ./output/template_analysis.json",
    )
    parser.add_argument(
        "--model",
        default=os.getenv("OPENAI_MODEL", "gpt-4.1"),
        help="LLM model name. Default: OPENAI_MODEL env or gpt-4.1",
    )
    parser.add_argument(
        "--mock",
        action="store_true",
        help="Skip API call and use built-in sample slide data.",
    )
    parser.add_argument(
        "--analyze-only",
        action="store_true",
        help="Only analyze the template and save template_analysis.json.",
    )
    parser.add_argument(
        "--skip-images",
        action="store_true",
        help="Skip image generation and keep the template image as-is. Overrides .env setting.",
    )
    parser.add_argument(
        "--google-safe",
        action="store_true",
        help="Generate a simpler PPTX for better Google Slides import compatibility.",
    )
    return parser.parse_args()


def ensure_dirs() -> None:
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    PROMPTS_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)


def find_default_template() -> Path:
    candidates = [
        path for path in TEMPLATE_DIR.glob("*.pptx") if not path.name.startswith("~$")
    ]
    if not candidates:
        raise FileNotFoundError(
            f"No .pptx template found in {TEMPLATE_DIR}. Put your template PPT there."
        )
    if len(candidates) == 1:
        return candidates[0]

    preferred = [path for path in candidates if path.stem != "templates"]
    pool = preferred or candidates
    return max(pool, key=lambda path: path.stat().st_mtime)


def build_output_path(template_path: Path) -> Path:
    return OUTPUT_DIR / f"generated_{template_path.stem}.pptx"


def build_json_output_path() -> Path:
    return OUTPUT_DIR / "generated_slide_plan.json"


def build_analysis_output_path() -> Path:
    return OUTPUT_DIR / "template_analysis.json"


def build_raw_output_path() -> Path:
    return OUTPUT_DIR / "llm_raw_response.txt"


def build_slide_plan_meta_path() -> Path:
    return OUTPUT_DIR / "generated_slide_plan.meta.json"


def build_session_plan_path(session_prefix: str) -> Path:
    return OUTPUT_DIR / f"{session_prefix}_slide_plan.json"


def build_session_raw_path(session_prefix: str) -> Path:
    return OUTPUT_DIR / f"{session_prefix}_llm_raw_response.txt"


def build_session_plan_meta_path(session_prefix: str) -> Path:
    return OUTPUT_DIR / f"{session_prefix}_slide_plan.meta.json"


def build_slide_plan_cache_key(prompt_text: str, model: str, use_mock: bool) -> str:
    payload = json.dumps(
        {
            "prompt_text": prompt_text,
            "model": model,
            "mock": use_mock,
        },
        ensure_ascii=False,
        sort_keys=True,
    )
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def load_cached_slide_plan(
    json_path: Path,
    raw_path: Path,
    cache_key: str | None = None,
    meta_path: Path | None = None,
) -> tuple[dict[str, Any] | None, str]:
    if not json_path.exists():
        return None, ""

    if cache_key:
        if meta_path is None or not meta_path.exists():
            return None, ""
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            return None, ""
        if meta.get("cache_key") != cache_key:
            return None, ""

    try:
        plan = json.loads(json_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return None, ""

    raw_text = raw_path.read_text(encoding="utf-8") if raw_path.exists() else ""
    return plan, raw_text


def write_slide_plan_cache_meta(meta_path: Path, cache_key: str) -> None:
    meta_path.write_text(
        json.dumps({"cache_key": cache_key}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def build_image_output_dir() -> Path:
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    return IMAGES_DIR


def build_merged_output_path() -> Path:
    return OUTPUT_DIR / "final_merged_curriculum.pptx"


def build_error_log_path() -> Path:
    return OUTPUT_DIR / "error.log"


def append_error_log(message: str) -> None:
    path = build_error_log_path()
    with path.open("a", encoding="utf-8") as fp:
        fp.write(message.rstrip() + "\n")


def write_session_error_file(session_prefix: str, text: str) -> Path:
    path = OUTPUT_DIR / f"{session_prefix}_error.txt"
    path.write_text(text, encoding="utf-8")
    return path


def write_merge_error_file(text: str) -> Path:
    path = OUTPUT_DIR / "merge_error.log"
    path.write_text(text, encoding="utf-8")
    return path


def read_prompt_file(path: Path) -> str:
    if not path.exists():
        raise FileNotFoundError(f"Prompt file not found: {path}")
    return path.read_text(encoding="utf-8")


def slugify_korean(text: str, max_len: int = 40) -> str:
    normalized = unicodedata.normalize("NFKC", text)
    normalized = re.sub(r"[^\w가-힣]+", "_", normalized, flags=re.UNICODE)
    normalized = re.sub(r"_+", "_", normalized).strip("_")
    return normalized[:max_len] or "session"


def parse_curriculum_file(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []

    lines = [line.strip() for line in path.read_text(encoding="utf-8-sig").splitlines()]
    header_pattern = re.compile(
        r"^\s*(?:#{1,6}\s*)?(?:[^\w\d가-힣]*\s*)?(\d+)교시\s*[\.:：]\s*(.+)$"
    )
    meta_labels = {
        "메모",
        "참고",
        "실습",
        "주의",
        "비고",
        "목표",
        "준비물",
        "과제",
        "핵심 구조",
        "AI 활용",
        "핵심 메시지",
    }
    core_labels = {
        "핵심 내용",
        "내용",
        "프로젝트 목적",
        "프로젝트 순서",
        "각 STEP",
        "최종 결과물",
    }
    label_names = "|".join(
        re.escape(label)
        for label in sorted(meta_labels | core_labels, key=len, reverse=True)
    )
    labeled_meta_pattern = re.compile(rf"^\s*({label_names})\s*[:：]\s*(.+)$")
    bracket_meta_pattern = re.compile(rf"^\s*\[({label_names})\]\s*(.+)$")

    sessions: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    pending_label: str | None = None
    current_section = "core"

    for line in lines:
        if not line:
            continue
        if set(line) <= {"─", "-", "_", " "}:
            continue

        line = re.sub(r"^\s*#{1,6}\s*", "", line).strip()
        line = re.sub(r"^\s*[*_]{1,3}(.+?)[*_]{1,3}\s*$", r"\1", line).strip()

        header_match = header_pattern.match(line)
        if header_match:
            if current:
                sessions.append(current)
            current = {
                "session_no": int(header_match.group(1)),
                "title": header_match.group(2).strip(),
                "core_points": [],
                "meta": {},
            }
            pending_label = None
            current_section = "core"
            continue

        if current is None:
            continue

        normalized_label = re.sub(r"^[^\w\d가-힣\[]+\s*", "", line).strip().rstrip(":：")
        if normalized_label in core_labels:
            current_section = "core"
            pending_label = None
            continue
        if normalized_label in meta_labels:
            current_section = "meta"
            pending_label = normalized_label
            continue

        if line.startswith("🟩 [") or line.startswith("🟦 ["):
            continue

        if line.startswith("👉"):
            value = line.replace("👉", "").strip()
            if pending_label and value:
                current["meta"].setdefault(pending_label, []).append(value)
            elif value:
                current["core_points"].append(value)
            continue

        if pending_label:
            current["meta"].setdefault(pending_label, []).append(line)
            if current_section != "meta":
                pending_label = None
            continue

        labeled_meta_match = labeled_meta_pattern.match(line)
        if labeled_meta_match:
            label = labeled_meta_match.group(1).strip()
            value = labeled_meta_match.group(2).strip()
            if label in core_labels and value:
                current["core_points"].append(value)
            elif value:
                current["meta"].setdefault(label, []).append(value)
            continue

        bracket_meta_match = bracket_meta_pattern.match(line)
        if bracket_meta_match:
            label = bracket_meta_match.group(1).strip()
            value = bracket_meta_match.group(2).strip()
            if label in core_labels and value:
                current["core_points"].append(value)
            elif value:
                current["meta"].setdefault(label, []).append(value)
            continue

        normalized_line = re.sub(r"^\s*[-•·●▪◦]\s*", "", line).strip()
        normalized_line = re.sub(r"^\s*\d+[\.\)]\s*", "", normalized_line).strip()
        if normalized_line:
            if current_section == "meta" and pending_label:
                current["meta"].setdefault(pending_label, []).append(normalized_line)
            else:
                current["core_points"].append(normalized_line)

    if current:
        sessions.append(current)

    return sessions


def parse_lecture_selection(raw: str | None) -> set[int]:
    if not raw:
        return set()

    selected: set[int] = set()
    for token in raw.split(","):
        value = token.strip()
        if not value:
            continue
        if not value.isdigit():
            raise ValueError(
                f"Invalid lecture selection '{value}'. Use numbers like 1 or 1,3."
            )
        selected.add(int(value))
    return selected


def parse_page_selection(raw: str | None) -> set[int]:
    if not raw:
        return set()

    selected: set[int] = set()
    for token in raw.split(","):
        value = token.strip()
        if not value:
            continue
        if not value.isdigit():
            raise ValueError(
                f"Invalid page selection '{value}'. Use numbers like 1 or 1,3."
            )
        page_no = int(value)
        if page_no <= 0:
            raise ValueError("Page numbers must be 1 or greater.")
        selected.add(page_no)
    return selected


def filter_slide_plan_pages(
    plan: dict[str, Any], selected_pages: set[int]
) -> dict[str, Any]:
    if not selected_pages:
        return plan

    filtered_slides = [
        slide
        for slide in plan["slides"]
        if int(slide.get("slide_number", 0)) in selected_pages
    ]
    if not filtered_slides:
        selected_text = ", ".join(str(number) for number in sorted(selected_pages))
        raise ValueError(f"No matching slide numbers found in plan for: {selected_text}")

    filtered_plan = dict(plan)
    filtered_plan["slides"] = filtered_slides
    return filtered_plan


def render_session_prompt(prompt_template: str, session: dict[str, Any]) -> str:
    core_lines = "\n".join(f"- {item}" for item in session["core_points"])
    meta_lines: list[str] = []
    for label, values in session.get("meta", {}).items():
        for value in values:
            meta_lines.append(f"- {label}: {value}")

    rendered = prompt_template.replace(
        "주제: 생성형 AI 개요 및 업무 변화 이해",
        f"주제: {session['title']}",
    )

    default_core_block = (
        "[핵심 내용]\n"
        "- 생성형 AI란 무엇인가 (LLM 개념)\n"
        "- 기존 자동화 vs 생성형 AI 차이\n"
        "- 기업에서의 활용 사례 (보고서, 분석, 요약)\n"
        "- 생성형 AI의 한계 (환각, 통제 불가)"
    )
    replacement_core_block = "[핵심 내용]\n" + core_lines
    if meta_lines:
        replacement_core_block += "\n[추가 맥락]\n" + "\n".join(meta_lines)

    rendered = rendered.replace(default_core_block, replacement_core_block)
    rendered += (
        f"\n\n[교시 정보]\n- 교시: {session['session_no']}교시\n- 이번 교시 주제: {session['title']}\n"
    )
    return rendered


def resolve_section_label(plan: dict[str, Any], lecture_title: str | None = None) -> str:
    if lecture_title:
        return lecture_title
    return plan.get("section_label", plan.get("deck_title", "강의안"))


def build_image_output_dir_for_session(session_prefix: str | None = None) -> Path:
    if not session_prefix:
        return build_image_output_dir()
    path = build_image_output_dir() / session_prefix
    path.mkdir(parents=True, exist_ok=True)
    return path


def format_page_suffix(selected_pages: set[int]) -> str:
    if not selected_pages:
        return ""
    joined = "_".join(str(number) for number in sorted(selected_pages))
    return f"_page{joined}"


def remove_all_shapes(slide) -> None:
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def clone_slide(prs: Presentation, source_index: int):
    source_slide = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    remove_all_shapes(new_slide)

    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    for rel in source_slide.part.rels.values():
        if any(
            token in rel.reltype
            for token in ("notesSlide", "slideLayout", "slideMaster", "theme")
        ):
            continue
        new_slide.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    return new_slide


def clone_slide_from_source(source_slide, target_prs: Presentation):
    raise NotImplementedError("Use clone_slide() within the template presentation instead.")


def remove_slides_by_indices(prs: Presentation, indices: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    for index in sorted(indices, reverse=True):
        slide_id = slide_id_list[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def get_picture_shapes(slide) -> list[Any]:
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def get_main_picture_spec(slide) -> dict[str, Any] | None:
    pictures = get_picture_shapes(slide)
    if not pictures:
        return None

    target = max(pictures, key=lambda shape: shape.width * shape.height)
    return {
        "left": target.left,
        "top": target.top,
        "width": target.width,
        "height": target.height,
    }


def remove_picture_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        shape.element.getparent().remove(shape.element)


def set_text(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.paragraphs[0].text = text


def normalize_paragraph_layout(paragraph, *, align=PP_ALIGN.LEFT) -> None:
    paragraph.alignment = align
    paragraph.left_indent = Pt(0)
    paragraph.first_line_indent = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def set_paragraphs(shape, lines: list[Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    cleaned = [line for line in lines if line]
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    if not cleaned:
        text_frame.paragraphs[0].text = ""
        return

    for idx, entry in enumerate(cleaned):
        if isinstance(entry, dict):
            line = str(entry.get("text") or "").strip()
            level = int(entry.get("level", 0) or 0)
            use_bullet = bool(entry.get("bullet", True))
        else:
            line = str(entry).strip()
            level = 0
            use_bullet = True

        if not line:
            continue

        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        normalize_paragraph_layout(paragraph, align=PP_ALIGN.LEFT)
        if GOOGLE_SAFE_MODE:
            indent = "    " * max(0, level)
            paragraph.level = 0
            paragraph.text = ""
            remove_bullet_from_paragraph(paragraph)
            prefix = f"{indent}• " if use_bullet else indent
            run = paragraph.add_run()
            if use_bullet:
                run.text = f"{prefix}{line}"
            else:
                run.text = f"{prefix}{line}"
        else:
            paragraph.text = line
            paragraph.level = max(0, level)
            if use_bullet:
                apply_bullet_to_paragraph(paragraph)
            else:
                remove_bullet_from_paragraph(paragraph)


def apply_bullet_to_paragraph(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()

    for node in list(p_pr):
        if any(
            node.tag.endswith(suffix)
            for suffix in ("}buNone", "}buAutoNum", "}buBlip", "}buChar")
        ):
            p_pr.remove(node)

    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    p_pr.append(bullet)


def remove_bullet_from_paragraph(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()

    for node in list(p_pr):
        if any(
            node.tag.endswith(suffix)
            for suffix in ("}buNone", "}buAutoNum", "}buBlip", "}buChar")
        ):
            p_pr.remove(node)

    bullet_none = OxmlElement("a:buNone")
    p_pr.append(bullet_none)


def apply_font_style(shape, font_name: str, font_size: int, bold: bool | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        paragraph.line_spacing = 1.5
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(34, 34, 34)
            if bold is not None:
                run.font.bold = bold


def apply_default_text_style(shape, font_name: str, font_size: int = 20) -> None:
    apply_font_style(shape, font_name=font_name, font_size=font_size, bold=None)


def apply_section_label_style(shape, font_name: str) -> None:
    apply_font_style(shape, font_name=font_name, font_size=20, bold=True)
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    for paragraph in shape.text_frame.paragraphs:
        normalize_paragraph_layout(paragraph, align=PP_ALIGN.CENTER)
        paragraph.line_spacing = 1.0


def apply_slide_title_style(shape, font_name: str) -> None:
    apply_font_style(shape, font_name=font_name, font_size=20, bold=True)
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    for paragraph in shape.text_frame.paragraphs:
        normalize_paragraph_layout(paragraph, align=PP_ALIGN.LEFT)
        paragraph.line_spacing = 1.0


def iter_text_shapes(slide) -> list[Any]:
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def table_shape_count(slide) -> int:
    return sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))


def picture_shape_count(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def shape_text_preview(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    return " | ".join(parts)[:120]


def suggest_tags(slide_number: int, text_shapes: list[Any], pictures: int, tables: int) -> list[str]:
    tags: list[str] = []

    if slide_number == 1:
        tags.append("cover")
    if slide_number == 2:
        tags.append("agenda")
    if slide_number in {3, 26}:
        tags.append("section_or_closing")
    if tables:
        tags.append("table")
    if pictures:
        tags.append("image")
    if len(text_shapes) >= 5:
        tags.append("multi_text")
    if len(text_shapes) >= 3 and not tables:
        tags.append("bullet")
    if len(text_shapes) >= 4 and not tables:
        tags.append("practice_like")
    return tags


def analyze_template(template_path: Path) -> dict[str, Any]:
    prs = Presentation(str(template_path))
    slides_summary: list[dict[str, Any]] = []
    candidate_templates: list[dict[str, Any]] = []
    content_template_slide = 2 if len(prs.slides) >= 2 else 1

    for slide_number, slide in enumerate(prs.slides, start=1):
        text_shapes = iter_text_shapes(slide)
        text_previews = [shape_text_preview(shape) for shape in text_shapes if shape_text_preview(shape)]
        pictures = picture_shape_count(slide)
        tables = table_shape_count(slide)
        tags = suggest_tags(slide_number, text_shapes, pictures, tables)

        slide_info = {
            "slide_number": slide_number,
            "text_shape_count": len(text_shapes),
            "picture_count": pictures,
            "table_count": tables,
            "tags": tags,
            "text_previews": text_previews[:5],
        }
        slides_summary.append(slide_info)

        if slide_number == content_template_slide:
            candidate_templates.append(
                {
                    "source_slide": slide_number,
                    "text_shape_count": len(text_shapes),
                    "picture_count": pictures,
                    "table_count": tables,
                    "tags": tags,
                    "use_when": "Use this as the main content slide template with left-top bullets and right-bottom image.",
                    "text_previews": text_previews[:4],
                }
            )

    return {
        "template_file": str(template_path),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "content_template_slide": content_template_slide,
        "candidate_templates": candidate_templates,
        "slides": slides_summary,
    }


def fill_table(shape, table_data: dict[str, Any]) -> None:
    if not getattr(shape, "has_table", False):
        return

    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    table = shape.table

    for col_idx, header in enumerate(headers[: len(table.columns)]):
        table.cell(0, col_idx).text = str(header)

    for row_idx, row in enumerate(rows[: len(table.rows) - 1], start=1):
        for col_idx, cell in enumerate(row[: len(table.columns)]):
            table.cell(row_idx, col_idx).text = str(cell)

    for row_idx, row in enumerate(table.rows):
        is_header = row_idx == 0
        for cell in row.cells:
            apply_table_cell_style(cell, is_header=is_header)


def set_cell_border(cell, side: str, color: str = "4C72B0", width: int = 12700) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    edge = tc_pr.find(qn(f"a:{side}"))
    if edge is None:
        edge = OxmlElement(f"a:{side}")
        tc_pr.append(edge)
    edge.set("w", str(width))
    edge.set("cap", "flat")
    edge.set("cmpd", "sng")
    edge.set("algn", "ctr")

    solid_fill = edge.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = OxmlElement("a:solidFill")
        edge.append(solid_fill)
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        srgb = OxmlElement("a:srgbClr")
        solid_fill.append(srgb)
    srgb.set("val", color)

    prst_dash = edge.find(qn("a:prstDash"))
    if prst_dash is None:
        prst_dash = OxmlElement("a:prstDash")
        edge.append(prst_dash)
    prst_dash.set("val", "solid")


def apply_table_cell_style(cell, is_header: bool = False) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = "맑은 고딕"
        paragraph.font.size = Pt(11)
        paragraph.line_spacing = 1.2
        paragraph.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "맑은 고딕"
            run.font.size = Pt(11)
            run.font.bold = is_header

    if not GOOGLE_SAFE_MODE:
        for side in ("lnL", "lnR", "lnT", "lnB"):
            set_cell_border(cell, side)


def get_first_table_spec(slide) -> dict[str, Any] | None:
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return {
                "rows": len(shape.table.rows),
                "cols": len(shape.table.columns),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
    return None


def add_blank_table_from_source(target_slide, source_slide, table_data: dict[str, Any]) -> None:
    table_spec = get_first_table_spec(source_slide)
    if not table_spec:
        fallback_spec = get_left_bottom_diagram_spec(source_slide)
        if not fallback_spec:
            return
        table_spec = {
            "rows": 4,
            "cols": 3,
            "left": fallback_spec["left"],
            "top": fallback_spec["top"],
            "width": fallback_spec["width"],
            "height": fallback_spec["height"],
        }

    rows = max(len(table_data.get("rows", [])) + 1, table_spec["rows"])
    cols = max(len(table_data.get("headers", [])), table_spec["cols"])
    remove_non_text_shapes_in_area(
        target_slide,
        {
            "left": table_spec["left"],
            "top": table_spec["top"],
            "width": table_spec["width"],
            "height": table_spec["height"],
        },
    )
    table_shape = target_slide.shapes.add_table(
        rows,
        cols,
        table_spec["left"],
        table_spec["top"],
        table_spec["width"],
        table_spec["height"],
    )
    fill_table(table_shape, table_data)


def remove_non_text_shapes_in_area(slide, area: dict[str, Any]) -> None:
    left = int(area["left"])
    top = int(area["top"])
    right = left + int(area["width"])
    bottom = top + int(area["height"])

    for shape in list(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            continue
        shape_left = int(shape.left)
        shape_top = int(shape.top)
        shape_right = shape_left + int(shape.width)
        shape_bottom = shape_top + int(shape.height)

        overlap = not (
            shape_right <= left
            or shape_left >= right
            or shape_bottom <= top
            or shape_top >= bottom
        )
        if overlap:
            shape.element.getparent().remove(shape.element)


def get_left_bottom_diagram_spec(source_slide) -> dict[str, Any] | None:
    picture_spec = get_main_picture_spec(source_slide)
    if not picture_spec:
        return None

    gap_x = max(int(picture_spec["width"] * 0.08), 120000)
    gap_y = max(int(picture_spec["height"] * 0.08), 120000)
    text_shapes = iter_text_shapes(source_slide)
    text_left = min((int(shape.left) for shape in text_shapes), default=0)
    text_bottom = max(
        (int(shape.top + shape.height) for shape in text_shapes),
        default=int(picture_spec["top"]),
    )

    left = text_left
    top = max(int(picture_spec["top"]), text_bottom + gap_y)
    right = int(picture_spec["left"] - gap_x)
    bottom = int(picture_spec["top"] + picture_spec["height"])
    width = right - left
    height = bottom - top

    min_width = max(int(picture_spec["width"] * 0.35), 1200000)
    min_height = max(int(picture_spec["height"] * 0.22), 900000)
    if width < min_width or height < min_height:
        return None

    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
    }


def build_diagram_positions(
    diagram_type: str,
    layout_direction: str,
    node_count: int,
    area: dict[str, Any],
) -> list[tuple[int, int, int, int]]:
    left = int(area["left"])
    top = int(area["top"])
    width = int(area["width"])
    height = int(area["height"])
    gap_x = max(int(width * 0.035), 100000)
    gap_y = max(int(height * 0.05), 90000)

    if node_count <= 0:
        return []

    if diagram_type in {"process", "comparison", "relationship"}:
        if layout_direction == "vertical":
            box_width = int(width * 0.70)
            box_height = int((height - gap_y * (node_count - 1)) / max(node_count, 1))
            x = left + int((width - box_width) / 2)
            return [
                (
                    x,
                    top + idx * (box_height + gap_y),
                    box_width,
                    box_height,
                )
                for idx in range(node_count)
            ]

        box_width = int((width - gap_x * (node_count - 1)) / max(node_count, 1))
        box_height = int(height * 0.58)
        y = top + int((height - box_height) / 2)
        return [
            (
                left + idx * (box_width + gap_x),
                y,
                box_width,
                box_height,
            )
            for idx in range(node_count)
        ]

    if diagram_type == "hierarchy":
        if node_count == 1:
            return [(left + int(width * 0.25), top + int(height * 0.2), int(width * 0.5), int(height * 0.3))]

        positions: list[tuple[int, int, int, int]] = []
        top_box_width = int(width * 0.42)
        top_box_height = int(height * 0.32)
        positions.append(
            (
                left + int((width - top_box_width) / 2),
                top + gap_y,
                top_box_width,
                top_box_height,
            )
        )

        bottom_count = node_count - 1
        bottom_gap = gap_x
        bottom_box_width = int((width - bottom_gap * (bottom_count - 1)) / max(bottom_count, 1))
        bottom_box_height = int(height * 0.30)
        bottom_y = top + height - bottom_box_height - int(gap_y * 0.6)
        for idx in range(bottom_count):
            positions.append(
                (
                    left + idx * (bottom_box_width + bottom_gap),
                    bottom_y,
                    bottom_box_width,
                    bottom_box_height,
                )
            )
        return positions

    if diagram_type == "cycle":
        if node_count == 1:
            return [(left + int(width * 0.25), top + int(height * 0.25), int(width * 0.5), int(height * 0.3))]

        box_width = int(width * 0.30)
        box_height = int(height * 0.24)
        center_x = left + width / 2
        center_y = top + height / 2
        radius_x = max((width - box_width) / 2.6, box_width / 2)
        radius_y = max((height - box_height) / 2.6, box_height / 2)
        positions = []
        for idx in range(node_count):
            angle = (2 * math.pi * idx / node_count) - (math.pi / 2)
            x = int(center_x + radius_x * math.cos(angle) - box_width / 2)
            y = int(center_y + radius_y * math.sin(angle) - box_height / 2)
            positions.append((x, y, box_width, box_height))
        return positions

    if layout_direction == "vertical":
        box_width = int(width * 0.70)
        box_height = int((height - gap_y * (node_count - 1)) / max(node_count, 1))
        x = left + int((width - box_width) / 2)
        return [
            (
                x,
                top + idx * (box_height + gap_y),
                box_width,
                box_height,
            )
            for idx in range(node_count)
        ]

    box_width = int((width - gap_x * (node_count - 1)) / max(node_count, 1))
    box_height = int(height * 0.58)
    y = top + int((height - box_height) / 2)
    return [
        (
            left + idx * (box_width + gap_x),
            y,
            box_width,
            box_height,
        )
        for idx in range(node_count)
    ]


def add_diagram_box(slide, text: str, left: int, top: int, width: int, height: int):
    return add_diagram_shape(
        slide,
        text=text,
        left=left,
        top=top,
        width=width,
        height=height,
        shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )


def add_diagram_shape(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    shape_type,
):
    shape = slide.shapes.add_shape(
        shape_type,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DIAGRAM_FILL_RGB
    shape.line.color.rgb = DIAGRAM_LINE_RGB
    shape.line.width = 1
    set_text(shape, text)
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    shape.text_frame.margin_left = Pt(7)
    shape.text_frame.margin_right = Pt(7)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    font_size = 12
    text_len = len(text.strip())
    if text_len >= 24:
        font_size = 9
    elif text_len >= 16:
        font_size = 10
    elif text_len >= 10:
        font_size = 11
    apply_default_text_style(shape, font_name="Malgun Gothic", font_size=font_size)
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.0
    return shape


def render_process_diagram(slide, positions, nodes) -> list[Any]:
    shapes = []
    for node_text, (left, top, width, height) in zip(nodes, positions):
        shapes.append(
            add_diagram_shape(
                slide,
                text=node_text,
                left=left,
                top=top,
                width=width,
                height=height,
                shape_type=MSO_AUTO_SHAPE_TYPE.CHEVRON,
            )
        )
    return shapes


def render_comparison_diagram(slide, positions, nodes) -> list[Any]:
    shapes = []
    for idx, (node_text, (left, top, width, height)) in enumerate(zip(nodes, positions)):
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.HEXAGON
            if idx == 0 or idx == len(nodes) - 1
            else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        )
        shapes.append(
            add_diagram_shape(
                slide,
                text=node_text,
                left=left,
                top=top,
                width=width,
                height=height,
                shape_type=shape_type,
            )
        )
    return shapes


def render_cycle_diagram(slide, positions, nodes) -> list[Any]:
    shapes = []
    for node_text, (left, top, width, height) in zip(nodes, positions):
        shapes.append(
            add_diagram_shape(
                slide,
                text=node_text,
                left=left,
                top=top,
                width=width,
                height=height,
                shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
            )
        )
    return shapes


def render_relationship_diagram(slide, positions, nodes) -> list[Any]:
    shapes = []
    center_index = len(nodes) // 2
    for idx, (node_text, (left, top, width, height)) in enumerate(zip(nodes, positions)):
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.OVAL
            if idx == center_index
            else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        )
        shapes.append(
            add_diagram_shape(
                slide,
                text=node_text,
                left=left,
                top=top,
                width=width,
                height=height,
                shape_type=shape_type,
            )
        )
    return shapes


def render_hierarchy_diagram(slide, positions, nodes) -> list[Any]:
    shapes = []
    for idx, (node_text, (left, top, width, height)) in enumerate(zip(nodes, positions)):
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.OVAL
            if idx == 0
            else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        )
        shapes.append(
            add_diagram_shape(
                slide,
                text=node_text,
                left=left,
                top=top,
                width=width,
                height=height,
                shape_type=shape_type,
            )
        )
    return shapes


def infer_diagram_type(slide_data: dict[str, Any], diagram_data: dict[str, Any]) -> str:
    raw_type = str(diagram_data.get("diagram_type") or "").strip().lower()
    if raw_type in {"process", "comparison", "hierarchy", "cycle", "relationship"}:
        if raw_type != "process":
            return raw_type

    title = str(slide_data.get("title") or "")
    bullets = " ".join(str(item) for item in (slide_data.get("bullets") or []))
    text = f"{title} {bullets}"

    if any(keyword in text for keyword in ["비교", "차이", "장단점", "대비"]):
        return "comparison"
    if any(keyword in text for keyword in ["순환", "반복", "루프", "사이클"]):
        return "cycle"
    if any(keyword in text for keyword in ["계층", "구조", "분류", "상위", "하위"]):
        return "hierarchy"
    if any(keyword in text for keyword in ["관계", "연결", "연동", "역할", "주체", "MCP", "API"]):
        return "relationship"
    return "process"


def infer_diagram_direction(diagram_type: str, slide_data: dict[str, Any], diagram_data: dict[str, Any]) -> str:
    raw_direction = str(diagram_data.get("layout_direction") or "").strip().lower()
    if raw_direction in {"horizontal", "vertical", "radial"}:
        return raw_direction

    if diagram_type == "cycle":
        return "radial"
    if diagram_type == "hierarchy":
        return "vertical"

    title = str(slide_data.get("title") or "")
    bullets = " ".join(
        str(item.get("text") if isinstance(item, dict) else item)
        for item in (slide_data.get("bullets") or [])
    )
    text = f"{title} {bullets}"
    if any(keyword in text for keyword in ["단계", "절차", "순서", "위에서 아래", "입력 후", "다음 단계"]):
        return "vertical"
    return "horizontal"


def add_connector_line(
    slide,
    from_box,
    to_box,
    diagram_type: str = "process",
    layout_direction: str = "horizontal",
) -> None:
    from_cx = int(from_box.left + from_box.width / 2)
    from_cy = int(from_box.top + from_box.height / 2)
    to_cx = int(to_box.left + to_box.width / 2)
    to_cy = int(to_box.top + to_box.height / 2)

    dx = to_cx - from_cx
    dy = to_cy - from_cy

    if layout_direction == "vertical" and diagram_type in {"process", "comparison", "relationship"}:
        x1 = from_cx
        y1 = int(from_box.top + from_box.height) if dy >= 0 else int(from_box.top)
        x2 = to_cx
        y2 = int(to_box.top) if dy >= 0 else int(to_box.top + to_box.height)
        connector_type = MSO_CONNECTOR.STRAIGHT
    elif diagram_type in {"process", "comparison", "relationship"}:
        x1 = int(from_box.left + from_box.width) if dx >= 0 else int(from_box.left)
        y1 = from_cy
        x2 = int(to_box.left) if dx >= 0 else int(to_box.left + to_box.width)
        y2 = to_cy
        connector_type = MSO_CONNECTOR.STRAIGHT
    elif diagram_type == "hierarchy":
        x1 = from_cx
        y1 = int(from_box.top + from_box.height)
        x2 = to_cx
        y2 = int(to_box.top)
        connector_type = MSO_CONNECTOR.STRAIGHT
    else:
        if abs(dx) >= abs(dy):
            x1 = int(from_box.left + from_box.width) if dx >= 0 else int(from_box.left)
            y1 = from_cy
            x2 = int(to_box.left) if dx >= 0 else int(to_box.left + to_box.width)
            y2 = to_cy
        else:
            x1 = from_cx
            y1 = int(from_box.top + from_box.height) if dy >= 0 else int(from_box.top)
            x2 = to_cx
            y2 = int(to_box.top) if dy >= 0 else int(to_box.top + to_box.height)
        connector_type = MSO_CONNECTOR.STRAIGHT

    connector = slide.shapes.add_connector(connector_type, x1, y1, x2, y2)
    connector.line.color.rgb = DIAGRAM_LINE_RGB
    connector.line.width = Pt(1.8 if diagram_type == "process" else 1.4)


def add_plain_line(
    slide, x1: int, y1: int, x2: int, y2: int, width_pt: float = 1.25, arrow_end: bool = False
):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = DIAGRAM_LINE_RGB
    connector.line.width = Pt(width_pt)
    return connector


def render_hierarchy_connectors(slide, node_shapes: list[Any]) -> None:
    if len(node_shapes) <= 1:
        return

    parent = node_shapes[0]
    children = node_shapes[1:]
    parent_cx = int(parent.left + parent.width / 2)
    parent_bottom = int(parent.top + parent.height)
    child_centers = [int(child.left + child.width / 2) for child in children]
    child_tops = [int(child.top) for child in children]
    vertical_gap = max(0, min(child_tops) - parent_bottom)
    trunk_y = int(parent_bottom + vertical_gap * 0.28)

    add_plain_line(slide, parent_cx, parent_bottom, parent_cx, trunk_y, width_pt=1.5, arrow_end=False)
    add_plain_line(
        slide,
        min(child_centers),
        trunk_y,
        max(child_centers),
        trunk_y,
        width_pt=1.5,
        arrow_end=False,
    )
    for child_cx, child_top in zip(child_centers, child_tops):
        add_plain_line(slide, child_cx, trunk_y, child_cx, child_top, width_pt=1.5, arrow_end=True)


def render_diagram(slide, source_slide, slide_data: dict[str, Any], diagram_data: dict[str, Any]) -> None:
    diagram_spec = get_left_bottom_diagram_spec(source_slide)
    if not diagram_spec:
        return

    nodes = [str(node).strip() for node in (diagram_data.get("nodes") or []) if str(node).strip()]
    if not nodes:
        return

    diagram_type = infer_diagram_type(slide_data, diagram_data)
    layout_direction = infer_diagram_direction(diagram_type, slide_data, diagram_data)
    remove_non_text_shapes_in_area(slide, diagram_spec)
    positions = build_diagram_positions(diagram_type, layout_direction, len(nodes), diagram_spec)
    if not positions:
        return

    if diagram_type == "process":
        node_shapes = render_process_diagram(slide, positions, nodes)
    elif diagram_type == "comparison":
        node_shapes = render_comparison_diagram(slide, positions, nodes)
    elif diagram_type == "cycle":
        node_shapes = render_cycle_diagram(slide, positions, nodes)
    elif diagram_type == "hierarchy":
        node_shapes = render_hierarchy_diagram(slide, positions, nodes)
    elif diagram_type == "relationship":
        node_shapes = render_relationship_diagram(slide, positions, nodes)
    else:
        node_shapes = []
        for node_text, (left, top, width, height) in zip(nodes, positions):
            node_shapes.append(add_diagram_box(slide, node_text, left, top, width, height))

    raw_links = diagram_data.get("links") or []
    if diagram_type == "hierarchy":
        render_hierarchy_connectors(slide, node_shapes)
    elif raw_links:
        link_map = {
            str(node).strip(): shape
            for node, shape in zip(nodes, node_shapes)
        }
        for link in raw_links:
            if not isinstance(link, (list, tuple)) or len(link) < 2:
                continue
            start_shape = link_map.get(str(link[0]).strip())
            end_shape = link_map.get(str(link[1]).strip())
            if start_shape and end_shape:
                add_connector_line(slide, start_shape, end_shape, diagram_type, layout_direction)
    else:
        for start_shape, end_shape in zip(node_shapes, node_shapes[1:]):
            add_connector_line(slide, start_shape, end_shape, diagram_type, layout_direction)

def normalize_bullet_line(text: str) -> str:
    cleaned = (text or "").strip()
    if not cleaned:
        return ""

    cleaned = re.sub(r"^[\u2022\u25E6\u25AA\u25CF\-\*\·\▪\‣]+\s*", "", cleaned)
    cleaned = re.sub(r"^\d+[\.\)]\s*", "", cleaned)
    cleaned = re.sub(r"^[가-힣A-Za-z][\.\)]\s*", "", cleaned)
    cleaned = re.sub(r"^[\u2022\u25E6\u25AA\u25CF\-\*\·\▪\‣]+\s*", "", cleaned)
    return cleaned.strip()


def parse_body_entry(item: Any) -> dict[str, Any] | None:
    if isinstance(item, dict):
        text = str(item.get("text") or "").strip()
        if not text:
            return None
        return {
            "text": normalize_bullet_line(text),
            "level": int(item.get("level", 0) or 0),
            "bullet": bool(item.get("bullet", True)),
        }

    raw = str(item or "").strip()
    if not raw:
        return None

    level = 0
    bullet = True
    text = raw

    if raw.lower().startswith("ex>"):
        return {
            "text": raw,
            "level": 0,
            "bullet": False,
        }

    hierarchy_match = re.match(r"^(>{1,3}|\-{1,3})\s*(.+)$", raw)
    if hierarchy_match:
        marker = hierarchy_match.group(1)
        text = hierarchy_match.group(2)
        level = max(0, len(marker) - 1)

    normalized = normalize_bullet_line(text)
    if not normalized:
        return None

    return {
        "text": normalized,
        "level": level,
        "bullet": bullet,
    }


def format_main_body(slide_data: dict[str, Any]) -> list[dict[str, Any]]:
    bullets = slide_data.get("bullets") or []
    normalized: list[dict[str, Any]] = []
    for bullet in bullets:
        entry = parse_body_entry(bullet)
        if entry:
            # Keep the body content as plain lines so PowerPoint does not render
            # literal bullet glyphs when template/style handling varies.
            entry["bullet"] = False
            normalized.append(entry)

    example_text = str(slide_data.get("example") or "").strip()
    if example_text:
        example_text = re.sub(r"^(ex>\s*)+", "", example_text, flags=re.IGNORECASE).strip()
        normalized.append(
            {
                "text": f"ex> {example_text}",
                "level": 0,
                "bullet": False,
            }
        )
    return normalized


def format_example_block(slide_data: dict[str, Any]) -> list[str]:
    lines: list[str] = []

    if slide_data.get("example"):
        lines.append("● 구체적 예시 또는 수치")
        lines.append(slide_data["example"])

    if slide_data.get("practice_prompt"):
        lines.append("● AI 프롬프트 예시")
        lines.append(slide_data["practice_prompt"])

    if slide_data.get("image_prompt"):
        lines.append("● 이미지 프롬프트")
        lines.append(slide_data["image_prompt"])

    return lines


def build_presenter_notes(slide_data: dict[str, Any]) -> list[str]:
    lines: list[str] = []

    if slide_data.get("why"):
        lines.append("[왜 이 내용이 필요한가]")
        lines.append(slide_data["why"])

    if slide_data.get("example"):
        lines.append("[구체적 예시 또는 수치]")
        lines.append(slide_data["example"])

    if slide_data.get("practice_prompt"):
        lines.append("[AI 프롬프트 예시]")
        lines.append(slide_data["practice_prompt"])

    if slide_data.get("image_prompt"):
        lines.append("[이미지 프롬프트]")
        lines.append(slide_data["image_prompt"])

    if slide_data.get("transition"):
        lines.append("[다음 단계 연결]")
        lines.append(slide_data["transition"])

    return lines


def set_presenter_notes(slide, slide_data: dict[str, Any], font_name: str) -> None:
    if GOOGLE_SAFE_MODE:
        return

    notes_slide = slide.notes_slide
    note_lines = build_presenter_notes(slide_data)
    if not note_lines:
        return

    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False):
            placeholder_type = shape.placeholder_format.type
            if str(placeholder_type) != "BODY (2)":
                continue
            set_paragraphs(shape, note_lines)
            apply_default_text_style(shape, font_name=font_name, font_size=12)
            return


def fill_slide(slide, source_slide, slide_data: dict[str, Any], section_label: str) -> None:
    font_name = "Malgun Gothic"
    slide_role = slide_data["slide_role"]
    text_shapes = iter_text_shapes(slide)

    if len(text_shapes) < 3:
        raise ValueError(
            f"Source slide {slide_data['source_slide']} does not have enough text slots."
        )

    set_text(text_shapes[0], section_label)
    set_text(text_shapes[1], slide_data["title"])
    apply_section_label_style(text_shapes[0], font_name=font_name)
    apply_slide_title_style(text_shapes[1], font_name=font_name)

    main_lines = format_main_body(slide_data)
    if slide_role in {"objective", "content", "summary"}:
        set_paragraphs(text_shapes[2], main_lines)
    elif slide_role == "practice_problem":
        set_paragraphs(text_shapes[2], main_lines)
    elif slide_role == "practice_answer":
        set_paragraphs(text_shapes[2], main_lines)
    else:
        set_paragraphs(text_shapes[2], main_lines)

    apply_default_text_style(text_shapes[2], font_name=font_name, font_size=16)
    set_presenter_notes(slide, slide_data, font_name=font_name)

    if slide_data.get("table"):
        add_blank_table_from_source(slide, source_slide, slide_data["table"])

    if slide_data.get("visual_type") == "diagram" and slide_data.get("diagram"):
        render_diagram(slide, source_slide, slide_data, slide_data["diagram"])


class LLMJSONParseError(RuntimeError):
    def __init__(self, message: str, raw_text: str):
        super().__init__(message)
        self.raw_text = raw_text


def extract_json_block(text: str) -> dict[str, Any]:
    def decode_or_raise(candidate: str, label: str) -> dict[str, Any]:
        try:
            return json.loads(candidate)
        except json.JSONDecodeError as exc:
            lines = candidate.splitlines()
            error_line = lines[exc.lineno - 1] if 0 < exc.lineno <= len(lines) else ""
            pointer = " " * max(exc.colno - 1, 0) + "^"
            message = (
                f"{label} JSON parse failed at line {exc.lineno}, column {exc.colno}: {exc.msg}\n"
                f"{error_line}\n{pointer}"
            )
            raise LLMJSONParseError(message, text) from exc

    try:
        return decode_or_raise(text, "full response")
    except LLMJSONParseError:
        pass

    fenced = re.search(r"```json\s*(\{.*\})\s*```", text, re.DOTALL)
    if fenced:
        return decode_or_raise(fenced.group(1), "fenced block")

    bare = re.search(r"(\{.*\})", text, re.DOTALL)
    if bare:
        return decode_or_raise(bare.group(1), "bare block")

    raise LLMJSONParseError("LLM response did not contain valid JSON.", text)


def build_llm_prompt(prompt_text: str, template_analysis: dict[str, Any]) -> str:
    analysis_payload = {
        "slide_count": template_analysis["slide_count"],
        "content_template_slide": template_analysis["content_template_slide"],
        "candidate_templates": template_analysis["candidate_templates"],
    }
    return (
        f"{prompt_text}\n\n"
        "[템플릿 분석 결과]\n"
        "아래 JSON은 템플릿 PPT를 로컬에서 분석한 결과다. "
        "슬라이드 구조를 설계할 때 반드시 이 후보 장표들 중에서 source_slide를 선택하라.\n\n"
        f"{json.dumps(analysis_payload, ensure_ascii=False, indent=2)}\n\n"
        f"{JSON_FORMAT_GUIDE}"
    )


def get_openai_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise EnvironmentError("OPENAI_API_KEY is not set.")

    return OpenAI(
        api_key=api_key,
        base_url=os.getenv("OPENAI_BASE_URL") or None,
        timeout=float(os.getenv("OPENAI_TIMEOUT_SECONDS", "300")),
    )


def generate_slide_plan_with_llm(
    prompt_text: str, template_analysis: dict[str, Any], model: str
) -> tuple[dict[str, Any], str]:
    client = get_openai_client()

    response = client.responses.create(
        model=model,
        input=[
            {
                "role": "system",
                "content": [
                    {
                        "type": "input_text",
                        "text": "You are a lecture deck planner. Output JSON only.",
                    }
                ],
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_text",
                        "text": build_llm_prompt(prompt_text, template_analysis),
                    }
                ],
            },
        ],
    )

    raw_text = response.output_text
    try:
        return extract_json_block(raw_text), raw_text
    except LLMJSONParseError as exc:
        raise LLMJSONParseError(str(exc), raw_text) from exc


def generate_slide_image(
    client: OpenAI,
    image_prompt: str,
    output_dir: Path,
    slide_number: int,
    image_size: str,
) -> Path:
    image_model = os.getenv("OPENAI_IMAGE_MODEL", "gpt-image-1")
    result = client.images.generate(
        model=image_model,
        prompt=enhance_image_prompt(image_prompt),
        size=image_size,
    )

    image_path = get_cached_image_path(output_dir, slide_number, image_prompt)
    image_base64 = None
    if getattr(result, "data", None):
        item = result.data[0]
        image_base64 = getattr(item, "b64_json", None)

    if not image_base64:
        raise ValueError("Image generation response did not include b64_json data.")

    image_path.write_bytes(base64.b64decode(image_base64))
    add_landscape_safe_margins(image_path)
    return image_path


def get_cached_image_path(
    output_dir: Path, slide_number: int, image_prompt: str | None = None
) -> Path:
    prompt_text = (image_prompt or "").strip()
    if prompt_text:
        prompt_hash = hashlib.sha256(prompt_text.encode("utf-8")).hexdigest()[:10]
        return output_dir / f"slide_{slide_number:02d}_{prompt_hash}_safe.png"
    return output_dir / f"slide_{slide_number:02d}_safe.png"


def enhance_image_prompt(image_prompt: str) -> str:
    base_prompt = (image_prompt or "").strip()
    if not base_prompt:
        return base_prompt

    safety_suffix = (
        " wide horizontal composition, landscape illustration, centered subject, "
        "leave generous empty padding on the left and right edges, keep all important subjects "
        "inside the central 70 percent of the image, no objects touching the left border, "
        "safe margins on all sides, no cropped head or hands, no important details near edges, "
        "no embedded text, no letters, no words, no labels"
    )
    return f"{base_prompt}, {safety_suffix}"


def add_landscape_safe_margins(image_path: Path) -> None:
    try:
        with Image.open(image_path) as img:
            original = img.convert("RGBA")
    except Exception:
        return

    width, height = original.size
    if width <= height or width <= 0 or height <= 0:
        return

    left_margin = int(round(width * LANDSCAPE_IMAGE_LEFT_SAFE_MARGIN_RATIO))
    right_margin = int(round(width * LANDSCAPE_IMAGE_SAFE_MARGIN_RATIO))
    vertical_margin = int(round(height * LANDSCAPE_IMAGE_SAFE_MARGIN_RATIO))
    content_width = max(1, width - left_margin - right_margin)
    content_height = max(1, height - (vertical_margin * 2))

    image_ratio = width / height
    box_ratio = content_width / content_height
    if image_ratio >= box_ratio:
        resized_width = content_width
        resized_height = max(1, int(round(content_width / image_ratio)))
    else:
        resized_height = content_height
        resized_width = max(1, int(round(content_height * image_ratio)))

    resized = original.resize((resized_width, resized_height), Image.Resampling.LANCZOS)
    canvas = Image.new("RGBA", (width, height), (255, 255, 255, 255))
    paste_left = left_margin + max(0, (content_width - resized_width) // 2)
    paste_top = vertical_margin + max(0, (content_height - resized_height) // 2)
    canvas.alpha_composite(resized, (paste_left, paste_top))
    canvas.convert("RGB").save(image_path)


def get_target_image_size_for_slide(source_slide) -> str:
    picture_spec = get_main_picture_spec(source_slide)
    if not picture_spec:
        return "1536x1024"

    width = int(picture_spec["width"])
    height = int(picture_spec["height"])
    if height <= 0:
        return "1536x1024"

    return "1536x1024" if (width / height) >= 1.0 else "1024x1536"


def image_matches_target_size(image_path: Path, expected_size: str) -> bool:
    try:
        expected_width, expected_height = [int(value) for value in expected_size.split("x", 1)]
        with Image.open(image_path) as img:
            width, height = img.size
        return width == expected_width and height == expected_height
    except Exception:
        return False


def fit_image_within_box(
    image_path: Path, box_left: int, box_top: int, box_width: int, box_height: int
) -> tuple[int, int, int, int]:
    with Image.open(image_path) as img:
        image_width_px, image_height_px = img.size

    if image_width_px <= 0 or image_height_px <= 0:
        return box_left, box_top, box_width, box_height

    box_ratio = box_width / box_height
    image_ratio = image_width_px / image_height_px

    if image_ratio >= box_ratio:
        fitted_width = box_width
        fitted_height = int(round(box_width / image_ratio))
    else:
        fitted_height = box_height
        fitted_width = int(round(box_height * image_ratio))

    fitted_left = box_left + max(0, (box_width - fitted_width) // 2)
    fitted_top = box_top + max(0, (box_height - fitted_height) // 2)
    return fitted_left, fitted_top, fitted_width, fitted_height


def replace_slide_image(slide, source_slide, image_path: Path) -> None:
    picture_spec = get_main_picture_spec(source_slide)
    if not picture_spec:
        return

    remove_picture_shapes(slide)
    left, top, width, height = fit_image_within_box(
        image_path,
        picture_spec["left"],
        picture_spec["top"],
        picture_spec["width"],
        picture_spec["height"],
    )
    slide.shapes.add_picture(
        str(image_path),
        left,
        top,
        width,
        height,
    )


def clone_rendered_slide_to_prs(source_slide, target_prs: Presentation):
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    for rel in source_slide.part.rels.values():
        if any(
            token in rel.reltype
            for token in ("notesSlide", "slideLayout", "slideMaster", "theme")
        ):
            continue
        new_slide.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    return new_slide


def decode_subprocess_output(data: bytes | None) -> str:
    if not data:
        return ""

    encodings = ["utf-8", "cp949"]
    preferred = locale.getpreferredencoding(False)
    if preferred and preferred not in encodings:
        encodings.append(preferred)

    for encoding in encodings:
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue

    return data.decode("utf-8", errors="replace")


def is_valid_presentation(path: Path) -> bool:
    if not path.exists():
        return False
    try:
        prs = Presentation(str(path))
        return len(prs.slides) > 0
    except Exception:
        return False


def merge_presentations(ppt_paths: list[Path], output_path: Path) -> None:
    if not ppt_paths:
        return
    merge_script = BASE_DIR / "merge_ppt.ps1"
    if not merge_script.exists():
        raise FileNotFoundError(f"Merge script not found: {merge_script}")

    merge_list_path = OUTPUT_DIR / "merge_input_files.txt"
    merge_list_path.write_text(
        "\n".join(str(ppt_path) for ppt_path in ppt_paths),
        encoding="utf-8",
    )

    command = [
        "powershell",
        "-NoProfile",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(merge_script),
        "-OutputFile",
        str(output_path),
        "-InputListFile",
        str(merge_list_path),
    ]

    result = subprocess.run(command, capture_output=True, text=False)
    stdout_text = decode_subprocess_output(result.stdout)
    stderr_text = decode_subprocess_output(result.stderr)

    if result.returncode != 0 and is_valid_presentation(output_path):
        warning_text = (
            f"PowerPoint merge returned a non-zero exit code but the merged file is valid.\n"
            f"Exit code: {result.returncode}\n"
            f"STDOUT:\n{stdout_text}\n"
            f"STDERR:\n{stderr_text}"
        )
        write_merge_error_file(warning_text)
        return

    if result.returncode != 0:
        error_text = (
            f"PowerPoint merge failed.\n"
            f"Exit code: {result.returncode}\n"
            f"STDOUT:\n{stdout_text}\n"
            f"STDERR:\n{stderr_text}"
        )
        write_merge_error_file(error_text)
        raise RuntimeError(error_text)


def extract_prompt_topic_and_core(prompt_text: str) -> tuple[str, list[str], list[str]]:
    topic = "강의 주제"
    topic_match = re.search(r"^주제:\s*(.+)$", prompt_text, re.MULTILINE)
    if topic_match:
        topic = topic_match.group(1).strip()

    core_points: list[str] = []
    core_match = re.search(
        r"\[핵심 내용\]\s*(.*?)(?=\n\[[^\]]+\]|\Z)",
        prompt_text,
        re.DOTALL,
    )
    if core_match:
        for line in core_match.group(1).splitlines():
            normalized = re.sub(r"^\s*[-•·●▪◦]\s*", "", line).strip()
            if normalized:
                core_points.append(normalized)

    context_points: list[str] = []
    context_match = re.search(
        r"\[추가 맥락\]\s*(.*?)(?=\n\[[^\]]+\]|\Z)",
        prompt_text,
        re.DOTALL,
    )
    if context_match:
        for line in context_match.group(1).splitlines():
            normalized = re.sub(r"^\s*[-•·●▪◦]\s*", "", line).strip()
            if normalized:
                context_points.append(normalized)

    return topic, core_points, context_points


def make_mock_slide_title(point: str, fallback: str) -> str:
    cleaned = re.sub(r"^[^:：]+[:：]\s*", "", point).strip()
    cleaned = re.sub(r"[“”\"']", "", cleaned)
    cleaned = re.split(r"[\.。]| vs | → ", cleaned)[0].strip()
    return cleaned[:28] or fallback


def build_prompt_based_mock_plan(prompt_text: str) -> dict[str, Any]:
    topic, core_points, context_points = extract_prompt_topic_and_core(prompt_text)
    source_points = core_points + context_points
    if not source_points:
        source_points = [topic]

    slides: list[dict[str, Any]] = [
        {
            "slide_number": 1,
            "slide_role": "objective",
            "source_slide": 2,
            "title": "학습 목표",
            "why": f"{topic}의 핵심 흐름을 먼저 정리해야 이후 실습과 적용이 쉬움.",
            "bullets": [
                f"{topic}의 기본 개념을 이해함",
                "핵심 구조와 판단 기준을 구분함",
                "실무 적용 포인트를 예시로 확인함",
                "마지막 실습으로 내용을 정리함",
            ],
            "example": "",
            "practice_prompt": "",
            "transition": "먼저 전체 개념과 구조를 확인함",
            "image_prompt": f"교육용 프레젠테이션 일러스트, {topic} 학습 목표를 정리한 장면, 깔끔한 인포그래픽 스타일",
            "visual_type": "bullets",
        }
    ]

    for slide_number in range(2, 18):
        point = source_points[(slide_number - 2) % len(source_points)]
        title = make_mock_slide_title(point, f"{topic} 핵심 {slide_number - 1}")
        slides.append(
            {
                "slide_number": slide_number,
                "slide_role": "content",
                "source_slide": 2,
                "title": title,
                "why": f"{topic}를 이해하려면 이 항목의 의미와 활용 방식을 구분해야 함.",
                "bullets": [
                    point,
                    f"{topic} 관점에서 핵심 의미를 정리함",
                    "실무 판단 기준과 연결해 해석함",
                    "다음 단계에서 적용할 질문으로 전환함",
                ],
                "example": f"{topic} 상황에서 이 항목을 기준으로 설명하거나 비교함",
                "practice_prompt": "",
                "transition": "이 내용을 바탕으로 다음 항목의 연결 구조를 확인함",
                "image_prompt": f"교육용 프레젠테이션 일러스트, {title} 내용을 설명하는 업무 장면, 깔끔한 벡터 스타일",
                "visual_type": "bullets",
            }
        )

    slides.extend(
        [
            {
                "slide_number": 18,
                "slide_role": "practice_problem",
                "source_slide": 2,
                "title": "실습 문제",
                "why": "학습 내용을 실제 업무 상황에 적용하는 단계가 필요함.",
                "bullets": [
                    f"{topic}와 관련된 업무 상황을 하나 선택함",
                    "핵심 문제와 판단 기준을 문장으로 정리함",
                    "필요한 데이터나 질문을 함께 작성함",
                ],
                "example": "",
                "practice_prompt": f"{topic} 주제에 맞는 업무 상황을 입력하고 핵심 문제, 판단 기준, 다음 질문을 도출하라.",
                "transition": "다음 장표에서 정답 구조 예시를 확인함",
                "image_prompt": f"교육용 프레젠테이션 일러스트, {topic} 실습 문제를 수행하는 장면, 깔끔한 벡터 스타일",
                "visual_type": "bullets",
            },
            {
                "slide_number": 19,
                "slide_role": "practice_answer",
                "source_slide": 2,
                "title": "실습 문제 정답 예시",
                "why": "정답 구조를 보면 실습 결과를 더 쉽게 점검할 수 있음.",
                "bullets": [
                    "문제 정의를 한 문장으로 먼저 제시함",
                    "판단 기준과 필요한 데이터를 분리함",
                    "다음 실행 질문을 구체적으로 작성함",
                ],
                "example": f"{topic} 상황에서 문제, 기준, 데이터, 실행 질문 순서로 정리함",
                "practice_prompt": "",
                "transition": "마지막으로 전체 핵심을 요약함",
                "image_prompt": f"교육용 프레젠테이션 일러스트, {topic} 정답 구조 예시를 정리한 장면, 깔끔한 벡터 스타일",
                "visual_type": "bullets",
            },
            {
                "slide_number": 20,
                "slide_role": "summary",
                "source_slide": 2,
                "title": "핵심 요약",
                "why": "마지막에 핵심 흐름을 다시 묶어야 학습 내용이 남음.",
                "bullets": (core_points[:4] or source_points[:4]) + [f"{topic}는 구조화된 질문과 판단 기준이 중요함"],
                "example": "",
                "practice_prompt": "",
                "transition": "",
                "image_prompt": f"교육용 프레젠테이션 일러스트, {topic} 핵심 요약 인포그래픽, 깔끔한 벡터 스타일",
                "visual_type": "bullets",
            },
        ]
    )

    return {
        "deck_title": topic,
        "section_label": topic,
        "slides": slides,
    }


def build_mock_plan() -> dict[str, Any]:
    slides: list[dict[str, Any]] = [
        {
            "slide_number": 1,
            "slide_role": "objective",
            "source_slide": 2,
            "title": "학습 목표",
            "why": "이번 강의에서 생성형 AI의 개념과 업무 변화의 핵심 흐름을 먼저 이해할 필요가 있음.",
            "bullets": [
                "생성형 AI와 LLM의 개념을 구분해 이해함",
                "기존 자동화와 생성형 AI 자동화의 차이를 파악함",
                "기업 업무에서 활용되는 대표 사례를 살펴봄",
                "환각과 통제 한계까지 함께 이해함",
            ],
            "example": "보고서 작성, 문서 요약, 데이터 분석 설명 같은 실무 장면을 기준으로 학습함",
            "practice_prompt": "",
            "transition": "먼저 생성형 AI가 무엇인지부터 이해해야 이후 활용 방식도 자연스럽게 연결됨",
            "image_prompt": "교육용 프레젠테이션 일러스트, 생성형 AI 학습 목표를 4개 키워드로 정리한 장면, 깔끔한 인포그래픽 스타일",
            "visual_type": "bullets",
        }
    ]

    content_slides: list[dict[str, Any]] = [
        {
            "title": "생성형 AI란 무엇인가",
            "why": "생성형 AI의 정의를 먼저 이해해야 이후 활용 방식과 한계를 정확히 설명할 수 있음.",
            "bullets": [
                "생성형 AI는 새 결과물을 만들어 내는 AI임",
                "기존 데이터를 분류하는 AI와 역할이 다름",
                "문장, 이미지, 요약처럼 새로운 출력을 생성함",
                "입력 방식에 따라 결과 품질이 크게 달라짐",
            ],
            "example": "회의 내용을 입력하면 보고서 초안이 바로 생성되는 방식임",
        },
        {
            "title": "LLM은 어떤 역할을 하는가",
            "why": "생성형 AI의 중심에 있는 언어 모델의 역할을 알아야 답변 생성 원리를 이해할 수 있음.",
            "bullets": [
                "LLM은 대량의 언어 데이터를 학습한 모델임",
                "입력 문맥을 보고 다음 표현을 예측함",
                "이 과정을 반복해 문장 단위의 답변을 만듦",
                "정답 검색보다 확률 기반 생성에 가까움",
            ],
            "example": "질문을 주면 가장 그럴듯한 문장 흐름으로 답을 이어 붙임",
        },
        {
            "title": "생성형 AI가 기존 AI와 다른 이유",
            "why": "기존 AI와의 차이를 알아야 생성형 AI의 업무 활용 포인트가 분명해짐.",
            "bullets": [
                "기존 AI는 분류와 예측 중심임",
                "생성형 AI는 문서와 결과물 생성 중심임",
                "비정형 업무에도 적용 범위가 넓음",
                "사용자 지시 방식이 결과를 좌우함",
            ],
            "example": "기존 모델은 스팸 메일을 분류하고, 생성형 AI는 답장 초안을 작성함",
        },
        {
            "title": "기존 자동화는 어떤 방식으로 동작했는가",
            "why": "생성형 AI 자동화와 비교하려면 기존 자동화의 구조부터 이해할 필요가 있음.",
            "bullets": [
                "기존 자동화는 규칙 기반으로 움직였음",
                "정해진 입력과 절차가 있어야 안정적으로 실행됨",
                "반복 업무에는 강하지만 예외 대응은 약함",
                "문서 해석이나 요약 같은 비정형 업무에는 한계가 큼",
            ],
            "example": "정해진 양식의 엑셀 파일을 읽어 시스템에 입력하는 방식임",
        },
        {
            "title": "생성형 AI 자동화는 무엇이 달라졌는가",
            "why": "생성형 AI 자동화의 변화 지점을 이해해야 업무 재설계 포인트가 보임.",
            "bullets": [
                {"text": "생성형 AI 자동화는 해석과 생성이 함께 들어감", "level": 0},
                {"text": "문서를 읽고 의미를 파악한 뒤 초안을 작성함", "level": 1},
                {"text": "사용자 요청에 따라 결과 표현을 바꿀 수 있음", "level": 1},
                {"text": "예외 상황에서도 설명 가능한 답을 제시함", "level": 1},
            ],
            "example": "주간 업무 보고 초안을 만들고 담당자가 검토해 최종본을 확정하는 구조임",
        },
        {
            "title": "두 방식의 차이를 표로 비교",
            "why": "기존 자동화와 생성형 AI 자동화를 한 화면에서 비교하면 차이가 빠르게 정리됨.",
            "bullets": [
                "업무 조건과 출력 방식의 차이를 비교해 볼 필요가 있음",
                "적용 가능한 업무 범위도 함께 비교하면 이해가 쉬움",
            ],
            "example": "정형 업무는 기존 자동화가 강하고, 비정형 문서 업무는 생성형 AI가 강함",
            "visual_type": "table",
            "table": {
                "headers": ["항목", "기존 자동화", "생성형 AI 자동화"],
                "rows": [
                    ["입력 조건", "정해진 형식 필요", "자연어 입력 가능"],
                    ["처리 방식", "규칙 순서대로 실행", "맥락 해석 후 생성"],
                    ["업무 범위", "반복 정형 업무", "문서 중심 비정형 업무"],
                    ["예외 대응", "설정 범위 안에서만 가능", "설명과 재작성까지 가능"],
                ],
            },
        },
        {
            "title": "보고서 작성 업무는 어떻게 바뀌는가",
            "why": "대표 업무 사례를 봐야 생성형 AI 도입 효과가 실제로 체감됨.",
            "bullets": [
                "핵심 메모만 있어도 초안 작성이 가능함",
                "문장 표현을 목적에 맞게 다시 정리할 수 있음",
                "보고 대상에 따라 톤과 구조를 바꿀 수 있음",
                "작성 시간보다 검토 시간이 더 중요해짐",
            ],
            "example": "회의 메모를 입력하면 팀장 보고용 문서 초안을 바로 생성함",
        },
        {
            "title": "데이터 분석 업무는 어떻게 바뀌는가",
            "why": "숫자 해석과 설명 작성에도 생성형 AI가 큰 역할을 하기 때문임.",
            "bullets": [
                "분석 결과를 설명 가능한 문장으로 바꿔 줌",
                "표와 차트의 핵심 의미를 먼저 정리해 줌",
                "초안 해석을 빠르게 만든 뒤 사람이 검증함",
                "결과 전달 속도가 크게 빨라질 수 있음",
            ],
            "example": "매출 증감 데이터를 입력하면 변화 원인을 설명하는 요약 문장을 제안함",
        },
        {
            "title": "문서 요약 업무는 어떻게 바뀌는가",
            "why": "긴 문서를 빠르게 이해하는 업무에서 생성형 AI의 효과가 분명하게 드러남.",
            "bullets": [
                "긴 문서를 핵심만 남겨 빠르게 요약함",
                "필요한 기준에 맞춰 요약 길이를 조절할 수 있음",
                "핵심 문장과 실행 항목을 따로 정리할 수 있음",
                "읽는 시간보다 판단 시간이 중요해짐",
            ],
            "example": "20페이지 보고서를 5개 핵심 메시지로 압축해 전달함",
        },
        {
            "title": "기업은 왜 생성형 AI를 도입하는가",
            "why": "도입 배경을 이해해야 현장에서 왜 관심이 큰지 설명할 수 있음.",
            "bullets": [
                "문서 업무 생산성을 빠르게 높일 수 있음",
                "반복되는 설명 작업을 줄일 수 있음",
                "의사결정용 초안 작성 속도를 올릴 수 있음",
                "기존 시스템을 크게 바꾸지 않아도 적용이 가능함",
            ],
            "example": "고객 문의 답변, 보고서 작성, 회의 요약처럼 바로 체감되는 업무부터 도입함",
        },
        {
            "title": "현장에서 자주 쓰는 활용 사례",
            "why": "실무 사례를 봐야 학습 내용이 업무 장면과 연결됨.",
            "bullets": [
                "보고서 초안 작성",
                "회의 내용 자동 요약",
                "분석 결과 설명 문장 생성",
                "이메일 답변 초안 작성",
            ],
            "example": "회의 녹취를 정리해 즉시 공유용 요약문으로 만드는 장면이 대표적임",
        },
        {
            "title": "활용 효과를 수치로 보면 무엇이 보이는가",
            "why": "도입 효과를 숫자로 보면 기대 수준과 적용 우선순위를 잡기 쉬움.",
            "bullets": [
                "초안 작성 시간 단축 효과가 가장 먼저 보임",
                "검토와 수정 시간을 포함해도 전체 소요가 줄어듦",
                "작은 업무부터 적용해도 체감 효과가 큼",
            ],
            "example": "주간 보고 작성 시간이 2시간에서 30분 수준으로 줄어드는 사례가 자주 나옴",
        },
        {
            "title": "생성형 AI의 가장 큰 한계는 무엇인가",
            "why": "효율만 강조하면 위험을 놓치기 때문에 한계도 함께 이해해야 함.",
            "bullets": [
                "그럴듯하지만 사실이 아닌 답을 만들 수 있음",
                "항상 같은 방식으로 통제되지 않음",
                "출력 품질이 질문 방식에 따라 크게 달라짐",
                "검토 없이 바로 쓰면 오류 위험이 커짐",
            ],
            "example": "존재하지 않는 통계 수치를 자연스럽게 문장 안에 넣는 경우가 있음",
        },
        {
            "title": "환각은 왜 발생하는가",
            "why": "환각의 원리를 알아야 생성형 AI 답변을 안전하게 검토할 수 있음.",
            "bullets": [
                {"text": "환각은 정답 검색이 아니라 확률 기반 생성에서 발생함", "level": 0},
                {"text": "그럴듯한 표현을 우선 이어 붙이는 성향이 있음", "level": 1},
                {"text": "출처 검증 없이 문장을 완성하려는 경향이 있음", "level": 1},
                {"text": "질문 맥락이 모호하면 오류 가능성이 더 커짐", "level": 1},
            ],
            "example": "없는 판례나 수치를 실제 사례처럼 말하는 장면이 대표적임",
        },
        {
            "title": "통제가 어려운 이유는 무엇인가",
            "why": "출력 통제의 어려움을 알아야 운영 기준과 검토 절차를 설계할 수 있음.",
            "bullets": [
                "같은 질문에도 표현과 결과가 달라질 수 있음",
                "질문 조건이 조금만 달라도 답변 방향이 바뀜",
                "모델 업데이트에 따라 출력 특성이 달라질 수 있음",
                "그래서 최종 판단은 사람의 검토가 필요함",
            ],
            "example": "같은 요약 요청이라도 강조하는 포인트가 매번 조금씩 달라질 수 있음",
        },
        {
            "title": "안전하게 쓰기 위한 운영 원칙",
            "why": "업무에 적용하려면 활용 기준과 검토 절차를 함께 갖춰야 함.",
            "bullets": [
                "초안 작성과 최종 승인 역할을 분리함",
                "출처 확인과 사실 검토 절차를 둠",
                "민감 정보 입력 기준을 명확히 정함",
                "결과물 품질 점검 기준을 함께 운영함",
            ],
            "example": "생성형 AI가 초안을 만들고 담당자가 근거와 표현을 검토하는 구조가 기본임",
        },
    ]

    for idx, slide in enumerate(content_slides, start=2):
        slides.append(
            {
                "slide_number": idx,
                "slide_role": "content",
                "source_slide": 2,
                "practice_prompt": "실무 상황을 기준으로 발표자가 바로 설명할 수 있는 질문 구조를 포함함",
                "transition": "이 내용을 바탕으로 다음 장표의 활용 흐름을 자연스럽게 이어감",
                "image_prompt": f"교육용 프레젠테이션 일러스트, {slide['title']} 내용을 설명하는 업무 장면, 깔끔한 벡터 스타일",
                "visual_type": slide.get("visual_type", "bullets"),
                **slide,
            }
        )

    slides.append(
        {
            "slide_number": 18,
            "slide_role": "practice_problem",
            "source_slide": 2,
            "title": "실습 문제",
            "why": "학습한 내용을 실제 업무 상황에 적용해 보는 단계가 필요함.",
            "bullets": [
                "팀 회의 내용을 보고용 문서로 바꾸는 문제를 설계함",
                "문제 정의, 요구사항, 출력 형식을 함께 제시함",
                "프롬프트를 직접 구성해 결과 구조를 확인함",
            ],
            "example": "한 주간 회의 메모를 바탕으로 팀장 보고용 1페이지 초안을 만드는 과제임",
            "practice_prompt": "회의 메모를 입력했을 때 요약, 핵심 이슈, 후속 액션을 포함한 초안이 나오도록 프롬프트를 설계함",
            "transition": "이제 정답 구조 예시를 보면서 좋은 프롬프트의 기준을 정리함",
            "image_prompt": "교육용 프레젠테이션 일러스트, 실습 문제를 수행하는 직장인과 노트북 화면, 깔끔한 벡터 스타일",
            "visual_type": "bullets",
        }
    )
    slides.append(
        {
            "slide_number": 19,
            "slide_role": "practice_answer",
            "source_slide": 2,
            "title": "실습 문제 정답 예시",
            "why": "좋은 프롬프트와 출력 구조의 예시를 보여 주면 실습 기준이 더 분명해짐.",
            "bullets": [
                "문제 정의, 요구사항, 출력 형식을 먼저 분리함",
                "역할과 대상 독자를 함께 지정함",
                "최종 결과 예시까지 포함해 품질 기준을 잡음",
            ],
            "example": "회의 결과를 5개 항목으로 정리하고, 후속 액션을 담당자 기준으로 분류하는 예시임",
            "practice_prompt": "당신은 팀 회의 내용을 정리하는 업무 보조자임. 아래 메모를 팀장 보고용으로 요약하고 핵심 이슈 3개와 후속 조치 2개를 제시하라.",
            "transition": "마지막으로 오늘 학습한 핵심 메시지를 다시 정리함",
            "image_prompt": "교육용 프레젠테이션 일러스트, 좋은 프롬프트 예시와 출력 결과를 비교하는 장면, 깔끔한 벡터 스타일",
            "visual_type": "bullets",
        }
    )
    slides.append(
        {
            "slide_number": 20,
            "slide_role": "summary",
            "source_slide": 2,
            "title": "핵심 요약",
            "why": "전체 흐름을 짧게 다시 정리해야 학습 내용이 구조적으로 남음.",
            "bullets": [
                "생성형 AI는 새로운 결과물을 만들어 내는 기술임",
                "LLM은 문맥을 보고 가장 그럴듯한 표현을 생성함",
                "기존 자동화와 달리 비정형 문서 업무까지 확장됨",
                "활용 효과가 크지만 환각과 통제 한계도 분명함",
                "안전한 활용을 위해서는 검토 절차가 반드시 필요함",
            ],
            "example": "",
            "practice_prompt": "",
            "transition": "",
            "image_prompt": "교육용 프레젠테이션 일러스트, 핵심 요약 5가지를 정리한 인포그래픽, 깔끔한 벡터 스타일",
            "visual_type": "bullets",
        }
    )

    mock_visual_overrides = {
        4: {
            "visual_type": "diagram",
            "diagram": {
                "diagram_type": "process",
                "title": "LLM 답변 생성 흐름",
                "nodes": ["질문 입력", "문맥 해석", "다음 표현 계산", "답변 생성"],
                "links": [
                    ["질문 입력", "문맥 해석"],
                    ["문맥 해석", "다음 표현 계산"],
                    ["다음 표현 계산", "답변 생성"],
                ],
                "notes": "입력부터 답변 생성까지의 흐름을 단계별로 설명함",
            },
        },
        7: {
            "visual_type": "diagram",
            "diagram": {
                "diagram_type": "comparison",
                "title": "기존 자동화와 생성형 AI 비교",
                "nodes": ["규칙 기반", "비교 포인트", "생성 기반"],
                "links": [
                    ["규칙 기반", "비교 포인트"],
                    ["비교 포인트", "생성 기반"],
                ],
                "notes": "두 방식의 입력 조건과 결과 생성 차이를 비교함",
            },
        },
        10: {
            "visual_type": "table",
            "table": {
                "headers": ["항목", "기존 방식", "생성형 AI"],
                "rows": [
                    ["초안 작성", "직접 작성", "초안 자동 생성"],
                    ["결과 설명", "수동 해석", "설명 문장 자동화"],
                    ["문서 요약", "전체 읽기", "핵심 요약 제공"],
                ],
            },
        },
        13: {
            "visual_type": "diagram",
            "diagram": {
                "diagram_type": "relationship",
                "title": "사람과 AI의 역할 분담",
                "nodes": ["업무 요청", "LLM", "초안 생성", "최종 검토"],
                "links": [
                    ["업무 요청", "LLM"],
                    ["LLM", "초안 생성"],
                    ["초안 생성", "최종 검토"],
                ],
                "notes": "AI는 초안을 만들고 사람은 판단과 검토를 담당함",
            },
        },
        16: {
            "visual_type": "diagram",
            "diagram": {
                "diagram_type": "cycle",
                "title": "안전한 활용 반복 구조",
                "nodes": ["질문 설계", "결과 생성", "사실 검토", "최종 반영"],
                "links": [
                    ["질문 설계", "결과 생성"],
                    ["결과 생성", "사실 검토"],
                    ["사실 검토", "최종 반영"],
                    ["최종 반영", "질문 설계"],
                ],
                "notes": "생성 후 검토를 거쳐 다시 개선하는 반복 구조를 보여 줌",
            },
        },
    }

    for slide in slides:
        slide.setdefault("visual_type", "bullets")
        override = mock_visual_overrides.get(slide["slide_number"])
        if override:
            slide.update(override)

    return {
        "deck_title": "생성형 AI 개요 및 업무 변화 이해",
        "section_label": "1. 생성형 AI와 프롬프트",
        "slides": slides,
    }


def normalize_slide_plan(plan: dict[str, Any], template_analysis: dict[str, Any]) -> dict[str, Any]:
    valid_source_slides = {
        item["source_slide"] for item in template_analysis["candidate_templates"]
    }
    fallback_source_slide = int(template_analysis.get("content_template_slide", 1))
    for slide in plan["slides"]:
        slide.setdefault("visual_type", "bullets")
        source_slide = slide.get("source_slide")
        if source_slide not in valid_source_slides:
            slide["source_slide"] = fallback_source_slide
    return plan


def generate_slide_plan(
    prompt_text: str,
    template_analysis: dict[str, Any],
    model: str,
    use_mock: bool,
) -> tuple[dict[str, Any], str]:
    if use_mock:
        plan = normalize_slide_plan(build_prompt_based_mock_plan(prompt_text), template_analysis)
        raw_text = json.dumps(plan, ensure_ascii=False, indent=2)
        return plan, raw_text
    plan, raw_text = generate_slide_plan_with_llm(prompt_text, template_analysis, model)
    return normalize_slide_plan(plan, template_analysis), raw_text


def render_presentation(
    template_path: Path,
    plan: dict[str, Any],
    output_path: Path,
    lecture_title: str | None = None,
) -> None:
    target_prs = Presentation(str(template_path))
    section_label = resolve_section_label(plan, lecture_title)

    original_slide_count = len(target_prs.slides)
    generated_slides: list[tuple[Any, Any, dict[str, Any]]] = []
    for slide_data in plan["slides"]:
        source_slide = target_prs.slides[slide_data["source_slide"] - 1]
        generated_slides.append(
            (clone_slide(target_prs, slide_data["source_slide"] - 1), source_slide, slide_data)
        )

    for slide, source_slide, slide_data in generated_slides:
        fill_slide(slide, source_slide, slide_data, section_label)

    remove_slides_by_indices(target_prs, list(range(original_slide_count)))
    target_prs.save(str(output_path))


def render_presentation_with_images(
    template_path: Path,
    plan: dict[str, Any],
    output_path: Path,
    skip_images: bool,
    image_output_dir: Path,
    lecture_title: str | None = None,
) -> None:
    target_prs = Presentation(str(template_path))
    section_label = resolve_section_label(plan, lecture_title)
    total_slides = len(plan["slides"])
    client = None if skip_images else get_openai_client()
    original_slide_count = len(target_prs.slides)

    for index, slide_data in enumerate(plan["slides"], start=1):
        print(f"[{index}/{total_slides}] 슬라이드 생성 중: {slide_data['title']}")
        source_slide = target_prs.slides[slide_data["source_slide"] - 1]
        slide = clone_slide(target_prs, slide_data["source_slide"] - 1)
        fill_slide(slide, source_slide, slide_data, section_label)
        if slide_data.get("visual_type") == "diagram" and slide_data.get("diagram"):
            print(f"[{index}/{total_slides}] 다이어그램 렌더링 완료")

        image_prompt = (slide_data.get("image_prompt") or "").strip()
        cached_image_path = get_cached_image_path(image_output_dir, index, image_prompt)
        target_image_size = get_target_image_size_for_slide(source_slide)

        if image_prompt and client and cached_image_path.exists():
            if client and not image_matches_target_size(cached_image_path, target_image_size):
                try:
                    cached_image_path.unlink()
                    print(f"[{index}/{total_slides}] 기존 정사각형 이미지 재생성 예정")
                except OSError:
                    pass
            else:
                replace_slide_image(slide, source_slide, cached_image_path)
                print(f"[{index}/{total_slides}] 기존 이미지 재사용")
                continue

        if image_prompt and client:
            print(f"[{index}/{total_slides}] 이미지 생성 중")
            try:
                image_path = generate_slide_image(
                    client=client,
                    image_prompt=image_prompt,
                    output_dir=image_output_dir,
                    slide_number=index,
                    image_size=target_image_size,
                )
                replace_slide_image(slide, source_slide, image_path)
                print(f"[{index}/{total_slides}] 이미지 적용 완료")
            except Exception as exc:
                print(f"[{index}/{total_slides}] 이미지 생성 실패, 템플릿 이미지 유지: {exc}")
        else:
            print(f"[{index}/{total_slides}] 이미지 생성 건너뜀")

    remove_slides_by_indices(target_prs, list(range(original_slide_count)))
    target_prs.save(str(output_path))




def main() -> None:
    global GOOGLE_SAFE_MODE
    load_env_file(DEFAULT_ENV_FILE)
    ensure_dirs()
    args = parse_args()
    GOOGLE_SAFE_MODE = True
    image_generation_enabled = env_flag("OPENAI_ENABLE_IMAGE_GENERATION", False)

    template_path = (
        Path(args.template).expanduser().resolve()
        if args.template
        else find_default_template()
    )
    curriculum_path = (
        Path(args.curriculum_file).expanduser().resolve()
        if args.curriculum_file
        else DEFAULT_CURRICULUM_FILE.resolve()
    )
    prompt_file = Path(args.prompt_file).expanduser().resolve()
    output_path = (
        Path(args.output).expanduser().resolve()
        if args.output
        else build_output_path(template_path)
    )
    json_output_path = (
        Path(args.json_output).expanduser().resolve()
        if args.json_output
        else build_json_output_path()
    )
    raw_output_path = build_raw_output_path()
    slide_plan_meta_path = build_slide_plan_meta_path()
    analysis_output_path = (
        Path(args.analysis_output).expanduser().resolve()
        if args.analysis_output
        else build_analysis_output_path()
    )

    template_analysis = analyze_template(template_path)
    analysis_output_path.write_text(
        json.dumps(template_analysis, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    if args.analyze_only:
        print(f"Template analysis saved: {analysis_output_path}")
        return

    prompt_template = read_prompt_file(prompt_file)
    sessions = parse_curriculum_file(curriculum_path)
    selected_lectures = parse_lecture_selection(args.lecture)
    selected_pages = parse_page_selection(args.page)

    if selected_pages and not selected_lectures:
        raise ValueError("--page can only be used together with --lecture.")

    if selected_lectures:
        sessions = [
            session for session in sessions if session["session_no"] in selected_lectures
        ]
        if not sessions:
            selected_text = ", ".join(str(number) for number in sorted(selected_lectures))
            raise ValueError(
                f"No matching lectures were found in {curriculum_path} for: {selected_text}"
            )

    if not sessions:
        cache_key = build_slide_plan_cache_key(prompt_template, args.model, args.mock)
        try:
            cached_plan, cached_raw_text = load_cached_slide_plan(
                json_output_path,
                raw_output_path,
                cache_key=cache_key,
                meta_path=slide_plan_meta_path,
            )
            if cached_plan is not None:
                print("[단일 작업] 기존 slide_plan 재사용")
                slide_plan = cached_plan
                raw_response_text = cached_raw_text
            else:
                print("[단일 작업] LLM 호출 시작")
                slide_plan, raw_response_text = generate_slide_plan(
                    prompt_template, template_analysis, args.model, args.mock
                )

            json_output_path.write_text(
                json.dumps(slide_plan, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            if raw_response_text:
                raw_output_path.write_text(raw_response_text, encoding="utf-8")
            write_slide_plan_cache_meta(slide_plan_meta_path, cache_key)
            render_presentation_with_images(
                template_path=template_path,
                plan=slide_plan,
                output_path=output_path,
                skip_images=args.mock or args.skip_images or (not image_generation_enabled),
                image_output_dir=build_image_output_dir_for_session(),
                lecture_title=slide_plan.get("deck_title"),
            )
        except Exception as exc:
            error_text = (
                "[단일 작업] 실패\n"
                f"오류 단계: LLM 호출 또는 PPT 생성\n"
                f"오류 타입: {type(exc).__name__}\n"
                f"오류 메시지: {exc}\n\n"
                f"{traceback.format_exc()}"
            )
            write_session_error_file("single_run", error_text)
            append_error_log(error_text)
            print("[단일 작업] 실패: output/single_run_error.txt 확인")
            raise

        print(f"Template analysis: {analysis_output_path}")
        print(f"Prompt file: {prompt_file}")
        print(f"LLM raw output: {raw_output_path}")
        print(f"Slide JSON: {json_output_path}")
        print(f"Created PPT: {output_path}")
        return

    total_sessions = len(sessions)
    print(f"Curriculum sessions detected: {total_sessions}")
    continue_on_error = env_flag("CONTINUE_ON_SESSION_ERROR", True)
    generated_ppt_paths: list[Path] = []
    page_suffix = format_page_suffix(selected_pages)

    for session_index, session in enumerate(sessions, start=1):
        session_slug = slugify_korean(session["title"])
        session_prefix = f"{session['session_no']:02d}_{session_slug}"
        if selected_pages:
            file_prefix = f"lecture{session['session_no']}{page_suffix}"
        else:
            file_prefix = session_prefix
        session_prompt_text = render_session_prompt(prompt_template, session)
        session_prompt_path = OUTPUT_DIR / f"{file_prefix}_prompt.txt"
        session_json_path = OUTPUT_DIR / f"{file_prefix}_slide_plan.json"
        session_raw_path = OUTPUT_DIR / f"{file_prefix}_llm_raw_response.txt"
        session_ppt_path = OUTPUT_DIR / f"{file_prefix}.pptx"
        full_session_json_path = build_session_plan_path(session_prefix)
        full_session_raw_path = build_session_raw_path(session_prefix)
        full_session_meta_path = build_session_plan_meta_path(session_prefix)
        cache_key = build_slide_plan_cache_key(session_prompt_text, args.model, args.mock)

        print(
            f"[교시 {session_index}/{total_sessions}] 생성 시작: {session['session_no']}교시 {session['title']}"
        )
        session_prompt_path.write_text(session_prompt_text, encoding="utf-8")

        try:
            cached_plan, cached_raw_text = load_cached_slide_plan(
                full_session_json_path,
                full_session_raw_path,
                cache_key=cache_key,
                meta_path=full_session_meta_path,
            )
            if cached_plan is not None:
                print(
                    f"[교시 {session_index}/{total_sessions}] 기존 slide_plan 재사용: {full_session_json_path}"
                )
                slide_plan = cached_plan
                raw_response_text = cached_raw_text
            else:
                print(f"[교시 {session_index}/{total_sessions}] LLM 호출 시작")
                slide_plan, raw_response_text = generate_slide_plan(
                    session_prompt_text, template_analysis, args.model, args.mock
                )

            full_slide_plan = slide_plan
            slide_plan = filter_slide_plan_pages(slide_plan, selected_pages)

            session_json_path.write_text(
                json.dumps(slide_plan, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            if raw_response_text:
                session_raw_path.write_text(raw_response_text, encoding="utf-8")
                full_session_raw_path.write_text(raw_response_text, encoding="utf-8")
            full_session_json_path.write_text(
                json.dumps(full_slide_plan, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            write_slide_plan_cache_meta(full_session_meta_path, cache_key)

            print(f"[교시 {session_index}/{total_sessions}] PPT 렌더링 시작")
            render_presentation_with_images(
                template_path=template_path,
                plan=slide_plan,
                output_path=session_ppt_path,
                skip_images=args.mock or args.skip_images or (not image_generation_enabled),
                image_output_dir=build_image_output_dir_for_session(file_prefix),
                lecture_title=session["title"],
            )

            generated_ppt_paths.append(session_ppt_path)
            print(f"[교시 {session_index}/{total_sessions}] 완료: {session_ppt_path}")
        except Exception as exc:
            raw_text = getattr(exc, "raw_text", None)
            if raw_text:
                session_raw_path.write_text(raw_text, encoding="utf-8")
            error_text = (
                f"[교시 {session_index}/{total_sessions}] 실패\n"
                f"교시: {session['session_no']}교시 {session['title']}\n"
                f"오류 타입: {type(exc).__name__}\n"
                f"오류 메시지: {exc}\n\n"
                f"{traceback.format_exc()}"
            )
            error_file = write_session_error_file(file_prefix, error_text)
            append_error_log(error_text)
            print(f"[교시 {session_index}/{total_sessions}] 실패: {error_file}")
            if not continue_on_error:
                raise

    merged_output_path = build_merged_output_path()
    if generated_ppt_paths and not selected_lectures:
        print(f"[병합] 교시별 PPT {len(generated_ppt_paths)}개를 하나로 병합 중")
        merge_presentations(generated_ppt_paths, merged_output_path)
        print(f"[병합] 완료: {merged_output_path}")

    print(f"Template analysis: {analysis_output_path}")
    print(f"Curriculum file: {curriculum_path}")
    print(f"Output folder: {OUTPUT_DIR}")


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("[취소] Ctrl+C로 생성 작업이 중단되었습니다.")
        raise SystemExit(130)
